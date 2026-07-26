# -*- coding: utf-8 -*-
"""발표.html → 발표.pptx (16:9, 편집 가능한 네이티브 PPT)

실행: python _발표-pptx빌드.py     (먼저 _발표-빌드.py로 발표.html을 만들어 둘 것)

흐름
 1) 발표.html을 1333.33x750px 슬라이드로 고정해 헤드리스 크롬에 띄우고
    `_발표-추출.js`가 요소별 위치·글자·색을 재서 로컬 수집 서버로 보낸다 (1px = 0.01in)
 2) 같은 HTML을 13.333x7.5in PDF로도 인쇄한다 — SVG 도해를 잘라낼 원본
 3) 잰 값으로 PowerPoint 도형·글상자·표를 만든다

편집 범위
 - 글자·도형·표: 전부 PowerPoint 개체라 그 자리에서 고칠 수 있다
 - 스크린샷과 SVG 도해: 이미지로 들어간다 (HTML에서 고치고 다시 돌려야 한다)

필요: Chrome 또는 Edge, PyMuPDF, python-pptx
"""
import io, os, re, json, base64, hashlib, sys, subprocess, threading, tempfile, shutil
from http.server import BaseHTTPRequestHandler, HTTPServer
import fitz
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.oxml.ns import qn

DECK = os.path.dirname(os.path.abspath(__file__))
SRC  = os.path.join(DECK, '발표.html')
OUT  = os.path.join(DECK, '발표.pptx')
PORT = 8940
BROWSERS = [r'C:\Program Files\Google\Chrome\Application\chrome.exe',
            r'C:\Program Files (x86)\Microsoft\Edge\Application\msedge.exe']

PX2IN = 0.01          # 1333.33px = 13.333in
PX2PT = 0.72
KR    = 'Malgun Gothic'
MONO  = 'Consolas'

TITLES = ["표지", "목차", "프로그램 UI 구성요소", "실행 순서", "작업 순서", "도면 4종",
          "예시 모델과 부재 분할", "제작도", "조립도", "설치도", "가공도", "값의 출처",
          "도면별 Osnap 선별 기준", "Osnap 선별 엔진", "동률일 때의 우선순위",
          "치수 보조선 배치", "제작도와 가공도의 배율 차이", "치수 선별 기준",
          "각도 표시 기준", "설치 위치 치수 기준", "부재 목록 수집", "도면 BOM 표 작성"]

MEASURE_CSS = """<style id="fix">
  html,body{ margin:0; padding:0; background:#fff; }
  .deck{ height:auto !important; overflow:visible !important; scroll-snap-type:none !important; }
  .page{ display:block !important; min-height:0 !important; padding:0 !important; margin:0 !important; }
  .slide{ width:1333.33px !important; height:750px !important; aspect-ratio:auto !important;
          border-radius:0 !important; box-shadow:none !important; }
  #bar,#hud{ display:none !important; }
</style>
"""

PRINT_CSS = """<style>
@media print{
  @page{ size:13.333in 7.5in; margin:0; }
  *{ -webkit-print-color-adjust:exact !important; print-color-adjust:exact !important; }
  html,body{ margin:0; padding:0; background:#fff; }
  .deck{ height:auto !important; overflow:visible !important; scroll-snap-type:none !important; }
  .page{ display:block !important; width:13.333in !important; height:7.5in !important;
         min-height:0 !important; padding:0 !important; margin:0 !important;
         page-break-after:always; break-after:page; page-break-inside:avoid; break-inside:avoid; }
  .page:last-child{ page-break-after:auto; break-after:auto; }
  .slide{ width:13.333in !important; height:7.5in !important; aspect-ratio:auto !important;
          border-radius:0 !important; box-shadow:none !important; }
  #bar,#hud{ display:none !important; }
}
</style>
"""

ASSET = None
_pdf = None


def browser():
    b = next((x for x in BROWSERS if os.path.exists(x)), None)
    if not b:
        sys.exit('Chrome/Edge를 찾지 못했습니다.')
    return b


# ================= 1) 요소 재기 =================
def measure(tmp):
    out = os.path.join(tmp, 'layout.json')
    done = threading.Event()

    class H(BaseHTTPRequestHandler):
        def do_POST(self):
            io.open(out, 'wb').write(self.rfile.read(int(self.headers.get('Content-Length', 0))))
            self.send_response(204); self.end_headers(); done.set()

        def log_message(self, *a):
            pass

    srv = HTTPServer(('127.0.0.1', PORT), H)
    threading.Thread(target=srv.serve_forever, daemon=True).start()
    try:
        html = io.open(SRC, encoding='utf-8').read()
        js = io.open(os.path.join(DECK, '_발표-추출.js'), encoding='utf-8').read()
        html = html.replace('</head>', MEASURE_CSS + '</head>', 1)
        html = html.replace('</body>', '<script>window.addEventListener("load",function(){'
                            'setTimeout(function(){' + js + '},300)});</script></body>', 1)
        page = os.path.join(tmp, 'measure.html')
        io.open(page, 'w', encoding='utf-8', newline='\n').write(html)
        subprocess.run([browser(), '--headless=new', '--disable-gpu', '--no-sandbox',
                        '--allow-file-access-from-files', '--window-size=1500,900',
                        '--virtual-time-budget=25000', '--dump-dom',
                        'file:///' + page.replace('\\', '/')], capture_output=True)
        if not done.wait(5) and not os.path.exists(out):
            sys.exit('요소 정보를 받지 못했습니다.')
        return json.load(io.open(out, encoding='utf-8'))
    finally:
        srv.shutdown()


# ================= 2) 도해 잘라낼 PDF =================
def render_pdf(tmp):
    html = io.open(SRC, encoding='utf-8').read().replace('</head>', PRINT_CSS + '</head>', 1)
    page = os.path.join(tmp, 'print.html')
    io.open(page, 'w', encoding='utf-8', newline='\n').write(html)
    pdf = os.path.join(tmp, 'deck.pdf')
    subprocess.run([browser(), '--headless=new', '--disable-gpu', '--no-sandbox',
                    '--no-pdf-header-footer', '--virtual-time-budget=20000',
                    '--print-to-pdf=' + pdf, 'file:///' + page.replace('\\', '/')],
                   check=True, capture_output=True)
    return pdf


# ================= 값 변환 =================
def rgb(c):
    """css color → (RGBColor, alpha 0~1). 없으면 (None, 0)"""
    if not c:
        return None, 0
    m = re.match(r'rgba?\(([\d.]+),\s*([\d.]+),\s*([\d.]+)(?:,\s*([\d.]+))?\)', c)
    if not m:
        return None, 0
    r, g, b = (int(round(float(m.group(i)))) for i in (1, 2, 3))
    a = float(m.group(4)) if m.group(4) is not None else 1.0
    return RGBColor(r, g, b), a


def IN(v):  return Inches(v * PX2IN)
def PTS(v): return Pt(max(1.0, v * PX2PT))


def align(a):
    return {'center': PP_ALIGN.CENTER, 'right': PP_ALIGN.RIGHT,
            'end': PP_ALIGN.RIGHT, 'justify': PP_ALIGN.JUSTIFY}.get(a, PP_ALIGN.LEFT)


def shape_kind(rect, rad):
    w, h = rect['w'], rect['h']
    # 정사각형에 가까울 때만 원. 알약 모양은 모서리를 최대로 둥글린 사각형으로
    if rad and rad >= min(w, h) * 0.48 and abs(w - h) <= min(w, h) * 0.25:
        return 'oval'
    if rad and rad >= 1.0:
        return 'round'
    return 'rect'


def set_fill(shp, bg):
    col, a = rgb(bg)
    if col is None or a == 0:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = col


def set_line(shp, borders):
    """네 변 중 공통 테두리만 도형 외곽선으로 준다. 한 변만 두꺼우면 목록으로 돌려준다."""
    if not borders:
        shp.line.fill.background(); return None
    base = min(b['w'] for b in borders)
    col, _ = rgb(borders[0]['c'])
    if base > 0 and col is not None:
        shp.line.color.rgb = col
        shp.line.width = Pt(max(0.5, base * PX2PT))
    else:
        shp.line.fill.background()
    return [(i, b) for i, b in enumerate(borders) if b['w'] > base + 0.6]


def accent_bars(slide, rect, thick):
    """border-left 4px 같은 강조 띠는 도형 외곽선으로 못 하니 얇은 사각형으로 따로 그린다."""
    for i, b in thick:
        col, _ = rgb(b['c'])
        if col is None:
            continue
        x, y, w, h = rect['x'], rect['y'], rect['w'], rect['h']
        box = [(x, y, w, b['w']), (x + w - b['w'], y, b['w'], h),
               (x, y + h - b['w'], w, b['w']), (x, y, b['w'], h)][i]
        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, IN(box[0]), IN(box[1]), IN(box[2]), IN(box[3]))
        s.fill.solid(); s.fill.fore_color.rgb = col
        s.line.fill.background(); s.shadow.inherit = False


def add_shape(slide, rect, rad, bg, borders):
    kind = shape_kind(rect, rad)
    mso = {'oval': MSO_SHAPE.OVAL, 'round': MSO_SHAPE.ROUNDED_RECTANGLE,
           'rect': MSO_SHAPE.RECTANGLE}[kind]
    shp = slide.shapes.add_shape(mso, IN(rect['x']), IN(rect['y']), IN(rect['w']), IN(rect['h']))
    if kind == 'round':
        shp.adjustments[0] = min(0.5, rad / min(rect['w'], rect['h']))
    set_fill(shp, bg)
    thick = set_line(shp, borders)
    shp.shadow.inherit = False
    if thick:
        accent_bars(slide, rect, thick)
    return shp


def fill_text(shp, item):
    tf = shp.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.auto_size = MSO_AUTO_SIZE.NONE
    lh = item.get('lh') or item['fs'] * 1.3
    single = item['rect']['h'] <= lh * 1.75
    # 브라우저에서 한 줄이던 것은 PowerPoint 글꼴 폭 차이로 접히지 않게 줄바꿈을 막는다
    tf.word_wrap = not single
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE if single else MSO_ANCHOR.TOP

    paras = [[]]
    for r in item['runs']:
        if r.get('br'):
            paras.append([])
        else:
            paras[-1].append(r)
    for idx, group in enumerate(paras):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = align(item.get('align'))
        p.line_spacing = Pt(lh * PX2PT)
        p.space_before = Pt(0); p.space_after = Pt(0)
        for r in group:
            run = p.add_run(); run.text = r['t']
            f = run.font
            f.size = PTS(r['fs'])
            f.name = MONO if r.get('mono') else KR
            f.bold = int(r.get('fw', 400) or 400) >= 600
            col, _ = rgb(r.get('color'))
            if col is not None:
                f.color.rgb = col
            rPr = run._r.get_or_add_rPr()
            if r.get('bg'):                       # code 칩 배경
                hc, _ = rgb(r['bg'])
                if hc is not None:
                    hl = rPr.makeelement(qn('a:highlight'), {})
                    hl.append(rPr.makeelement(qn('a:srgbClr'),
                                              {'val': '%02X%02X%02X' % (hc[0], hc[1], hc[2])}))
                    rPr.append(hl)
            # 동아시아 글꼴을 같이 지정해야 PowerPoint가 한글에 다른 글꼴을 끼워넣지 않는다
            for tag in ('a:ea', 'a:cs'):
                rPr.append(rPr.makeelement(qn(tag), {'typeface': MONO if r.get('mono') else KR}))


# ================= 이미지 =================
_imgcache = {}


def data_uri_to_file(src):
    key = hashlib.md5(src[:200].encode() + str(len(src)).encode()).hexdigest()[:12]
    if key in _imgcache:
        return _imgcache[key]
    fp = os.path.join(ASSET, key + '.png')
    io.open(fp, 'wb').write(base64.b64decode(src.split(',', 1)[1]))
    _imgcache[key] = fp
    return fp


def painted_rect(item):
    """object-fit: contain 인 이미지가 실제로 그려지는 영역"""
    r = item['rect']
    if item.get('fit') != 'contain' or not item.get('nw'):
        return r
    s = min(r['w'] / item['nw'], r['h'] / item['nh'])
    w, h = item['nw'] * s, item['nh'] * s
    return {'x': r['x'] + (r['w'] - w) / 2, 'y': r['y'] + (r['h'] - h) / 2, 'w': w, 'h': h}


def svg_crop(page_no, rect):
    """SVG 도해는 PDF에서 그 영역만 3배 크기로 잘라낸다."""
    fp = os.path.join(ASSET, 'svg%02d_%d_%d.png' % (page_no, int(rect['x']), int(rect['y'])))
    if not os.path.exists(fp):
        scale = 3.0
        m = fitz.Matrix(1333.33 / 960 * scale, 750 / 540 * scale)
        clip = fitz.Rect(rect['x'] / 1333.33 * 960, rect['y'] / 750 * 540,
                         (rect['x'] + rect['w']) / 1333.33 * 960,
                         (rect['y'] + rect['h']) / 750 * 540)
        _pdf[page_no - 1].get_pixmap(matrix=m, clip=clip, alpha=False).save(fp)
    return fp


# ================= 표 =================
def cell_bottom_border(cell, bb):
    """python-pptx가 노출하지 않는 칸 아래 테두리를 직접 넣는다."""
    if not bb or bb['w'] <= 0:
        return
    col, a = rgb(bb['c'])
    if col is None or a == 0:
        return
    tcPr = cell._tc.get_or_add_tcPr()
    ln = tcPr.makeelement(qn('a:lnB'), {'w': str(int(max(0.5, bb['w'] * PX2PT) * 12700)),
                                        'cap': 'flat', 'cmpd': 'sng', 'algn': 'ctr'})
    fill = tcPr.makeelement(qn('a:solidFill'), {})
    fill.append(tcPr.makeelement(qn('a:srgbClr'), {'val': '%02X%02X%02X' % (col[0], col[1], col[2])}))
    ln.append(fill)
    tcPr.append(ln)


def add_table(slide, data):
    rows = data['rows']
    ncol = max(sum(c['span'] for c in r['cells']) for r in rows)
    tr = data['rect']
    tbl = slide.shapes.add_table(len(rows), ncol, IN(tr['x']), IN(tr['y']),
                                 IN(tr['w']), IN(tr['h'])).table
    tbl.first_row = False
    tbl.horz_banding = False
    ref = max(rows, key=lambda r: len(r['cells']))
    if len(ref['cells']) == ncol:
        for i, c in enumerate(ref['cells']):
            tbl.columns[i].width = IN(c['rect']['w'])
    for ri, r in enumerate(rows):
        tbl.rows[ri].height = IN(r['rect']['h'])
        ci = 0
        for c in r['cells']:
            cell = tbl.cell(ri, ci)
            if c['span'] > 1:
                cell.merge(tbl.cell(ri, min(ncol - 1, ci + c['span'] - 1)))
            cell.margin_left = cell.margin_right = Inches(0.04)
            cell.margin_top = cell.margin_bottom = Inches(0.01)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            col, _ = rgb(c['bg']) if c['bg'] else (None, 0)
            if col is not None:
                cell.fill.solid(); cell.fill.fore_color.rgb = col
            else:
                cell.fill.background()
            fill_text(cell, {'runs': c['runs'], 'align': c['align'], 'fs': c['fs'],
                             'lh': c['fs'] * 1.35, 'rect': {'h': c['rect']['h']}})
            cell_bottom_border(cell, c['bb'])
            ci += c['span']


# ================= 3) 조립 =================
def build(lay, pdf, tmp):
    global ASSET, _pdf
    ASSET = os.path.join(tmp, 'assets'); os.makedirs(ASSET, exist_ok=True)
    _pdf = fitz.open(pdf)

    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    blank = prs.slide_layouts[6]
    for s in lay['slides']:
        slide = prs.slides.add_slide(blank)
        for it in s['items']:
            t = it['type']
            if t == 'box':
                add_shape(slide, it['rect'], it.get('rad', 0), it.get('bg'), it.get('borders'))
            elif t == 'text':
                if it.get('bg') or it.get('borders'):
                    shp = add_shape(slide, it['rect'], it.get('rad', 0), it.get('bg'), it.get('borders'))
                else:
                    r = it['rect']
                    shp = slide.shapes.add_textbox(IN(r['x']), IN(r['y']), IN(r['w']), IN(r['h']))
                fill_text(shp, it)
            elif t == 'img':
                pr = painted_rect(it)
                slide.shapes.add_picture(data_uri_to_file(it['src']),
                                         IN(pr['x']), IN(pr['y']), IN(pr['w']), IN(pr['h']))
            elif t == 'svg':
                r = it['rect']
                slide.shapes.add_picture(svg_crop(s['n'], r),
                                         IN(r['x']), IN(r['y']), IN(r['w']), IN(r['h']))
            elif t == 'table':
                add_table(slide, it['data'])
        slide.notes_slide.notes_text_frame.text = TITLES[s['n'] - 1]
    prs.save(OUT)
    print('=> %s (%d슬라이드, %.1f MB)' % (OUT, len(lay['slides']), os.path.getsize(OUT) / 1048576))


def main():
    tmp = tempfile.mkdtemp(prefix='deck-pptx-')
    try:
        lay = measure(tmp)
        if len(lay['slides']) != len(TITLES):
            sys.exit('슬라이드 수(%d)와 제목 수(%d)가 다릅니다 — TITLES를 갱신하세요.'
                     % (len(lay['slides']), len(TITLES)))
        build(lay, render_pdf(tmp), tmp)
    finally:
        shutil.rmtree(tmp, ignore_errors=True)


if __name__ == '__main__':
    main()
